from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
from PIL import Image
import os


def create_slide(header, content, image_url):
    # Create a presentation object
    presentation = Presentation()

    # Slide layout with title, content, and picture placeholders
    slide_layout = presentation.slide_layouts[5]  # Title Slide layout with content and picture placeholders

    # Add a slide
    slide = presentation.slides.add_slide(slide_layout)

    # Add the header
    title = slide.shapes.title
    title.text = header

    # Add the content
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    p = text_frame.add_paragraph()
    p.text = content

    # Add the image
    image_url = image_url
    response = requests.get(image_url)
    image_bytes = BytesIO(response.content)
    left = Inches(4)  # Adjust the position as needed
    top = Inches(1.5)
    width = Inches(4)  # Adjust the size as needed
    height = Inches(3)
    picture = slide.shapes.add_picture(image_bytes, left, top, width, height)

    return presentation


# Sample data
slide_data = [
    {
        "header": "Slide 1 Header hello hello hello hello hello hello hello hello helllo hello hello hello hello",
        "content": "Slide 1 Content",
        "image_url": "https://www.simplilearn.com/ice9/free_resources_article_thumb/what_is_image_Processing.jpg"
    },
    {
        "header": "Slide 2 Header",
        "content": "Slide 2 Content Slide 1 Header hello hello hello hello hello hello hello hello helllo hello hello hello hello",
        "image_url": 'https://www.simplilearn.com/ice9/free_resources_article_thumb/what_is_image_Processing.jpg'
    }
    # Add more slides if needed
]

# Create the slides
presentation = create_slide(slide_data[0]["header"], slide_data[0]["content"], slide_data[0]["image_url"])
for data in slide_data[:]:
    presentation = create_slide(data["header"], data["content"], data["image_url"])

# Save the presentation
presentation.save("slide_with_image.pptx")

# Open the presentation
os.startfile("slide_with_image.pptx")